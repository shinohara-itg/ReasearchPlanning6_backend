import os
import zipfile
from pathlib import Path

import fitz  # PyMuPDF
import openpyxl
from docx import Document
from pptx import Presentation


def safe_extract_zip(zip_path: str, extract_to: str) -> None:
    extract_base = Path(extract_to).resolve()

    with zipfile.ZipFile(zip_path, "r") as z:
        for member in z.infolist():
            dest = (extract_base / member.filename).resolve()
            if not str(dest).startswith(str(extract_base)):
                raise ValueError(f"Unsafe zip entry detected: {member.filename}")
        z.extractall(extract_to)


def read_txt(path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp932"):
        try:
            with open(path, "r", encoding=enc, errors="ignore") as f:
                return f.read()
        except Exception:
            continue
    return ""


def read_pdf(path: str) -> str:
    try:
        doc = fitz.open(path)
        return "\n".join(page.get_text("text") for page in doc)
    except Exception:
        return ""


def read_pptx_text(path: str) -> str:
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text") and shp.text:
                    texts.append(shp.text)
        return "\n".join(texts)
    except Exception:
        return ""


def read_docx(path: str) -> str:
    try:
        doc = Document(path)
        parts = []

        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)

        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [(c.text or "").replace("\n", " ").strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))

        return "\n".join(parts).strip()
    except Exception:
        return ""


def read_xlsx(path: str) -> str:
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        parts = []

        for ws in wb.worksheets:
            parts.append(f"【Sheet】{ws.title}")
            for row in ws.iter_rows(values_only=True):
                vals = []
                for v in row:
                    if v is None:
                        continue
                    s = str(v).strip()
                    if s:
                        vals.append(s)
                if vals:
                    parts.append(" | ".join(vals))

        return "\n".join(parts).strip()
    except Exception:
        return ""


def extract_text_from_file(path: str) -> str:
    low = path.lower()

    if low.endswith(".txt"):
        return read_txt(path)
    if low.endswith(".pdf"):
        return read_pdf(path)
    if low.endswith(".pptx"):
        return read_pptx_text(path)
    if low.endswith(".docx"):
        return read_docx(path)
    if low.endswith(".xlsx"):
        return read_xlsx(path)

    return ""


def extract_documents_from_path(path: str) -> list:
    documents = []
    low = path.lower()

    if low.endswith(".zip"):
        extract_dir = str(Path(path).parent / f"{Path(path).stem}_extracted")
        os.makedirs(extract_dir, exist_ok=True)

        try:
            safe_extract_zip(path, extract_dir)
        except Exception:
            return []

        for root, _, files in os.walk(extract_dir):
            for filename in files:
                file_path = os.path.join(root, filename)
                text = extract_text_from_file(file_path)
                if text.strip():
                    documents.append({
                        "file_name": filename,
                        "text": text
                    })

        return documents

    text = extract_text_from_file(path)
    if text.strip():
        documents.append({
            "file_name": Path(path).name,
            "text": text
        })

    return documents