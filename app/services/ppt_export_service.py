from __future__ import annotations

from io import BytesIO
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from app.schemas.ppt_export import PptShapeItem, PptTableItem

from pptx.dml.color import RGBColor
from pptx.util import Pt

def fill_schedule_tables(prs, payload):
    schedule_tables = payload.get("schedule_tables") or []
    if not schedule_tables:
        return

    slide = prs.slides[15]  # P16

    table_names = ["table1", "table2", "table3"]

    for table_index, shape_name in enumerate(table_names):
        if table_index >= len(schedule_tables):
            continue

        target_shape = None
        for shape in slide.shapes:
            if shape.name == shape_name:
                target_shape = shape
                break

        if target_shape is None or not target_shape.has_table:
            continue

        table = target_shape.table
        rows = schedule_tables[table_index]

        for i, row in enumerate(rows):
            table_row_index = i + 1  # 0行目はヘッダー

            if table_row_index >= len(table.rows):
                break

            date_text = row[0] if len(row) > 0 else ""
            client_text = row[1] if len(row) > 1 else ""
            itg_text = row[2] if len(row) > 2 else ""
            is_holiday = bool(row[3]) if len(row) > 3 else False

            values = [date_text, client_text, itg_text]

            for col_index, value in enumerate(values):
                cell = table.cell(table_row_index, col_index)
                cell.text = str(value or "")

                # 土日祝グレーアウト
                if is_holiday:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(230, 230, 230)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)


                text_frame = cell.text_frame
                text_frame.word_wrap = True

                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(6)
                    paragraph.font.name = "Yu Gothic"

                    if not paragraph.runs:
                        run = paragraph.add_run()
                        run.text = ""

                    for run in paragraph.runs:
                        run.font.size = Pt(6)
                        run.font.name = "Yu Gothic"

def _iter_shapes_recursive(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_recursive(shape.shapes)


def _apply_text_style(target, font_name: str = "Arial", font_size: int = 10) -> None:
    try:
        text_frame = getattr(target, "text_frame", None)
        if text_frame is None:
            return

        text_frame.word_wrap = True

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT

            # runがない場合に備えてparagraph.fontにも指定
            paragraph.font.name = font_name
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)

            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(0, 0, 0)

    except Exception:
        pass


def _set_text_to_named_shape(slide, shape_name: str, text: str) -> bool:
    for shape in _iter_shapes_recursive(slide.shapes):
        if getattr(shape, "name", "") != shape_name:
            continue

        if getattr(shape, "has_text_frame", False):
            shape.text = text
            _apply_text_style(shape)
            return True

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell.text = text
                        _apply_text_style(cell)
                return True
            except Exception:
                return False

        return False

    return False

def _set_table_to_named_shape(
    slide,
    shape_name: str,
    rows: list[list],
    start_row: int = 1,
) -> bool:
    for shape in _iter_shapes_recursive(slide.shapes):
        if getattr(shape, "name", "") != shape_name:
            continue

        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            return False

        try:
            table = shape.table

            max_rows = len(table.rows)
            max_cols = len(table.columns)

            # データ行クリア
            for r in range(start_row, max_rows):
                for c in range(max_cols):
                    table.cell(r, c).text = ""

            # データ投入
            for i, row_values in enumerate(rows):
                target_row = start_row + i

                if target_row >= max_rows:
                    break

                for c in range(max_cols):
                    value = row_values[c] if c < len(row_values) else ""

                    cell = table.cell(target_row, c)
                    cell.text = "" if value is None else str(value)

                    _apply_text_style(cell)

            return True

        except Exception:
            return False

    return False

def export_ppt_from_template_bytes(
    template_bytes: bytes,
    items: Iterable[PptShapeItem],
    tables: Iterable[PptTableItem] | None = None,
    payload: dict | None = None,
) -> tuple[BytesIO, dict]:
    prs = Presentation(BytesIO(template_bytes))

    report = {
        "total_items": 0,
        "applied": 0,
        "skipped_empty": 0,
        "slide_oob": 0,
        "shape_not_found": 0,
        "table_applied": 0,
        "errors": [],
    }

    for item in items:
        report["total_items"] += 1

        try:
            text = item.text or ""
            if not text.strip():
                report["skipped_empty"] += 1
                continue

            if item.slide_index < 0 or item.slide_index >= len(prs.slides):
                report["slide_oob"] += 1
                continue

            slide = prs.slides[item.slide_index]
            ok = _set_text_to_named_shape(
                slide=slide,
                shape_name=item.shape_name,
                text=text,
            )

            if ok:
                report["applied"] += 1
            else:
                report["shape_not_found"] += 1

        except Exception as e:
            report["errors"].append(
                {
                    "slide_index": item.slide_index,
                    "shape_name": item.shape_name,
                    "error": str(e),
                }
            )

    # -----------------------------
    # table items
    # -----------------------------
    if tables:
        for table_item in tables:
            try:
                if (
                    table_item.slide_index < 0
                    or table_item.slide_index >= len(prs.slides)
                ):
                    report["slide_oob"] += 1
                    continue

                slide = prs.slides[table_item.slide_index]

                ok = _set_table_to_named_shape(
                    slide=slide,
                    shape_name=table_item.shape_name,
                    rows=table_item.rows,
                    start_row=table_item.start_row,
                )

                if ok:
                    report["table_applied"] += 1
                else:
                    report["shape_not_found"] += 1

            except Exception as e:
                report["errors"].append(
                    {
                        "slide_index": table_item.slide_index,
                        "shape_name": table_item.shape_name,
                        "error": str(e),
                    }
                )

    if payload:
        fill_schedule_tables(prs, payload)

    output = BytesIO()
    prs.save(output)
    output.seek(0)

    return output, report